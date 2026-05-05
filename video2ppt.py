#!/usr/bin/env python3
"""video2ppt.py
CLI 工具：
 - 指定视频路径，按固定或可调时间间隔抽帧
 - 感知哈希去重（默认阈值 2）
 - 生成 16:9 填满裁切的 PPT，每张图片占一页

依赖：moviepy, pillow, imagehash, python-pptx, tqdm
用法示例：
    python video2ppt.py -i /path/to/video.mp4 --interval 2 --threshold 2
"""
import argparse
import itertools
import os
from pathlib import Path
import sys
import tempfile
from datetime import timedelta

from PIL import Image
import imagehash
try:
    from moviepy.editor import VideoFileClip  # MoviePy < 2.1
except ModuleNotFoundError:
    # 新版 MoviePy 2.1+ 将 editor 拆分
    from moviepy.video.io.VideoFileClip import VideoFileClip
from pptx import Presentation
from pptx.util import Inches, Emu
from tqdm import tqdm

# EMU 常量，pptx 里 1 inch = 914400 EMU
EMU_PER_INCH = 914400

# -------- PyInstaller 兼容处理 --------
# 如果脚本被 PyInstaller 打包为 exe，sys.frozen = True，
# 将打包目录加入 PATH，确保同目录下的 ffmpeg.exe 可被调用。
if getattr(sys, 'frozen', False):
    bundle_dir = Path(sys.executable).parent
    os.environ['PATH'] = str(bundle_dir) + os.pathsep + os.environ.get('PATH', '')


def seconds_to_timestamp(seconds: float) -> str:
    """把秒数格式化为 HH-MM-SS."""
    td = timedelta(seconds=int(seconds))
    # timedelta 默认到天，这里保证 HH 范围足够
    total_seconds = int(td.total_seconds())
    hours = total_seconds // 3600
    minutes = (total_seconds % 3600) // 60
    secs = total_seconds % 60
    return f"{hours:02d}-{minutes:02d}-{secs:02d}"


def extract_frames(video_path: Path, raw_dir: Path, interval: float):
    """按照 interval(s) 抽帧保存到 raw_dir。"""
    raw_dir.mkdir(parents=True, exist_ok=True)
    clip = VideoFileClip(str(video_path))
    duration = clip.duration
    timestamps = list(itertools.takewhile(lambda t: t < duration, itertools.count(0, interval)))

    for t in tqdm(timestamps, desc="Extracting frames", unit="frame"):
        frame = clip.get_frame(t)  # ndarray (H, W, 3)
        img = Image.fromarray(frame)
        name = f"{seconds_to_timestamp(t)}.jpg"
        img.save(raw_dir / name, "JPEG", quality=95)
    clip.reader.close()
    if clip.audio is not None:
        # 兼容不同版本 MoviePy
        if hasattr(clip.audio, "reader") and clip.audio.reader is not None:
            reader = clip.audio.reader
            if hasattr(reader, "close_proc"):
                reader.close_proc()
            elif hasattr(reader, "close"):
                reader.close()


def deduplicate_frames(raw_dir: Path, unique_dir: Path, threshold: int):
    """根据感知哈希去重，将唯一帧复制/移动到 unique_dir。"""
    unique_dir.mkdir(parents=True, exist_ok=True)
    hashes = []  # List[imagehash.ImageHash]

    files = sorted(raw_dir.glob("*.jpg"))
    for img_path in tqdm(files, desc="Deduplicating", unit="img"):
        img = Image.open(img_path)
        h = imagehash.phash(img)
        is_dup = False
        for old_h in hashes:
            if h - old_h <= threshold:  # 汉明距离
                is_dup = True
                break
        if not is_dup:
            hashes.append(h)
            # 保存唯一图片
            target_path = unique_dir / img_path.name
            img.save(target_path)


def add_full_bleed_picture(slide, img_path: Path, slide_width: int, slide_height: int):
    """在给定 slide 中插入填满页面的图片 (16:9)，保持比例并可能裁切。"""
    with Image.open(img_path) as im:
        img_width_px, img_height_px = im.size

    # 计算需填满的缩放因子 (EMU per px)
    scale = max(slide_width / img_width_px, slide_height / img_height_px)
    new_width = int(img_width_px * scale)
    new_height = int(img_height_px * scale)

    # 居中摆放，可能出现负 offset 让其溢出裁切
    left = int((slide_width - new_width) / 2)
    top = int((slide_height - new_height) / 2)

    slide.shapes.add_picture(str(img_path), Emu(left), Emu(top), Emu(new_width), Emu(new_height))


def create_ppt(unique_dir: Path, ppt_path: Path):
    """用 unique_dir 下的图片生成 PPT。"""
    prs = Presentation()
    # 设置 16:9 尺寸 13.33" x 7.5"
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # 空白布局

    slide_w, slide_h = prs.slide_width, prs.slide_height

    images = sorted(unique_dir.glob("*.jpg"))
    for img_path in tqdm(images, desc="Building PPT", unit="slide"):
        slide = prs.slides.add_slide(blank_layout)
        add_full_bleed_picture(slide, img_path, slide_w, slide_h)

    prs.save(str(ppt_path))


def parse_args():
    p = argparse.ArgumentParser(description="Extract frames from video, deduplicate, and export to PPT.")
    p.add_argument("-i", "--input", required=True, help="Path to input video file.")
    p.add_argument("--interval", type=float, default=2.0, help="Frame extraction interval in seconds (default: 2).")
    p.add_argument("--threshold", type=int, default=2, help="Perceptual hash Hamming distance threshold (default: 2).")
    return p.parse_args()


def main():
    args = parse_args()
    video_path = Path(args.input).expanduser().resolve()
    if not video_path.exists():
        sys.exit(f"Input video not found: {video_path}")

    # 输出目录：视频名 + _output
    output_dir = video_path.with_name(f"{video_path.stem}_output")
    frames_raw_dir = output_dir / "frames_raw"
    frames_unique_dir = output_dir / "frames_unique"
    ppt_path = output_dir / f"{video_path.stem}_unique.pptx"

    print(f"[1/3] Extracting frames to {frames_raw_dir} ...")
    extract_frames(video_path, frames_raw_dir, args.interval)

    print(f"[2/3] Deduplicating frames → {frames_unique_dir} (threshold={args.threshold}) ...")
    deduplicate_frames(frames_raw_dir, frames_unique_dir, args.threshold)

    print(f"[3/3] Building PPT {ppt_path} ...")
    create_ppt(frames_unique_dir, ppt_path)

    print("Done.")


if __name__ == "__main__":
    main() 